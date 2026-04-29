"""
Generate PPTX previews and PNG thumbnail images for all seed templates.

Pipeline per template:
  1. AI (Gemini) generates professional slide content as JSON
  2. build_premium_pptx() converts JSON → .pptx  (saved in seeds/previews/)
  3. PowerPoint COM exports slide 1 → high-res PNG  (saved in seeds/previews/images/)
  4. Pillow creates a web-ready 800×450 thumbnail  (seeds/previews/thumbnails/)
  5. DB record updated: preview_pptx_path + thumbnail_url

Usage:
    cd backend
    python seeds/generate_preview_images.py
    python seeds/generate_preview_images.py --force   # regenerate all
    python seeds/generate_preview_images.py --template "Executive Strategy"
"""
from __future__ import annotations

import argparse
import asyncio
import json
import os
import sys
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.config import settings
from app.core.database import close_db, init_db
import app.core.database as _db_module
from app.utils.logger import get_logger

logger = get_logger("preview_images")

SEEDS_DIR = Path(__file__).parent
THEMES_DIR = SEEDS_DIR / "themes"
TEMPLATES_DIR = SEEDS_DIR / "templates"
PREVIEWS_DIR = SEEDS_DIR / "previews"
IMAGES_DIR = PREVIEWS_DIR / "images"       # full-size slide exports
THUMBNAILS_DIR = PREVIEWS_DIR / "thumbnails"  # 800x450 web thumbnails

THUMB_W, THUMB_H = 800, 450


# ---------------------------------------------------------------------------
# PPTX → PNG via PowerPoint COM (Windows)
# ---------------------------------------------------------------------------

def pptx_slide_to_png(pptx_path: str, out_path: str, slide_index: int = 0,
                      width_px: int = 1920, height_px: int = 1080) -> bool:
    """Export a single PPTX slide to PNG using PowerPoint COM automation."""
    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()
    powerpoint = None
    prs = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1

        abs_pptx = str(Path(pptx_path).resolve())
        abs_out = str(Path(out_path).resolve())

        prs = powerpoint.Presentations.Open(abs_pptx, ReadOnly=True, Untitled=True, WithWindow=False)
        slide = prs.Slides(slide_index + 1)  # COM is 1-indexed
        slide.Export(abs_out, "PNG", width_px, height_px)
        logger.info(f"Exported slide → {abs_out}")
        return True
    except Exception as exc:
        logger.info(f"PowerPoint COM export failed: {exc}")
        return False
    finally:
        try:
            if prs:
                prs.Close()
        except Exception:
            pass
        try:
            if powerpoint:
                powerpoint.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()


# ---------------------------------------------------------------------------
# Pillow fallback: draw a styled thumbnail directly from theme + slide data
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def draw_thumbnail_from_theme(
    template_name: str,
    category: str,
    description: str,
    theme_colors: dict,
    out_path: str,
    width: int = THUMB_W,
    height: int = THUMB_H,
) -> None:
    """Draw a branded thumbnail card using Pillow when COM export is unavailable."""
    from PIL import Image, ImageDraw, ImageFont

    primary = _hex_to_rgb(theme_colors.get("primary", "#0A2342"))
    accent = _hex_to_rgb(theme_colors.get("accent", "#C0932F"))
    bg = _hex_to_rgb(theme_colors.get("background", "#FAFAFA"))
    text_col = _hex_to_rgb(theme_colors.get("text", "#1C2B3A"))

    img = Image.new("RGB", (width, height), bg)
    draw = ImageDraw.Draw(img)

    # Background gradient: left dark panel, right lighter
    panel_w = int(width * 0.55)
    for x in range(panel_w):
        ratio = x / panel_w
        r = int(primary[0] * (1 - ratio * 0.3))
        g = int(primary[1] * (1 - ratio * 0.3))
        b = int(primary[2] * (1 - ratio * 0.3))
        draw.line([(x, 0), (x, height)], fill=(r, g, b))

    # Right panel subtle tint
    draw.rectangle([panel_w, 0, width, height], fill=bg)

    # Accent bar
    bar_h = max(6, height // 80)
    bar_y = height // 2 - height // 8
    draw.rectangle([40, bar_y, 40 + 80, bar_y + bar_h], fill=accent)

    # Template name (title)
    font_size_title = max(32, width // 22)
    try:
        font_title = ImageFont.truetype("georgiab.ttf", font_size_title)
    except Exception:
        try:
            font_title = ImageFont.truetype("arial.ttf", font_size_title)
        except Exception:
            font_title = ImageFont.load_default()

    font_size_body = max(16, width // 50)
    try:
        font_body = ImageFont.truetype("arial.ttf", font_size_body)
    except Exception:
        font_body = ImageFont.load_default()

    font_size_cat = max(13, width // 70)
    try:
        font_cat = ImageFont.truetype("arialbd.ttf", font_size_cat)
    except Exception:
        font_cat = font_body

    # Category label
    draw.text((40, bar_y - font_size_cat - 20), category.upper(),
              font=font_cat, fill=accent)

    # Title text — wrap manually at ~28 chars
    title_words = template_name.split()
    lines, current = [], ""
    for word in title_words:
        candidate = f"{current} {word}".strip()
        if len(candidate) > 20:
            lines.append(current)
            current = word
        else:
            current = candidate
    if current:
        lines.append(current)

    ty = bar_y + bar_h + 18
    for line in lines[:3]:
        draw.text((40, ty), line, font=font_title, fill=(255, 255, 255))
        ty += font_size_title + 8

    # Description snippet
    desc_snippet = description[:80] + ("…" if len(description) > 80 else "")
    draw.text((40, ty + 12), desc_snippet, font=font_body, fill=(180, 200, 220))

    # Decorative right-side pattern: stacked accent lines
    rx = panel_w + 40
    for i in range(5):
        lw = int(width * 0.25 * (1 - i * 0.15))
        ly = int(height * 0.25) + i * int(height * 0.12)
        alpha = 255 - i * 40
        r2, g2, b2 = accent
        draw.rectangle([rx, ly, rx + lw, ly + 6],
                       fill=(r2, g2, b2, alpha) if img.mode == "RGBA" else (r2, g2, b2))

    # WACDeckStudio watermark bottom-right
    try:
        font_wm = ImageFont.truetype("arial.ttf", max(11, width // 80))
    except Exception:
        font_wm = ImageFont.load_default()
    wm_text = "WACDeckStudio"
    bbox = draw.textbbox((0, 0), wm_text, font=font_wm)
    wm_w = bbox[2] - bbox[0]
    draw.text((width - wm_w - 20, height - 28), wm_text,
              font=font_wm, fill=(120, 140, 160))

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path, "PNG", optimize=True)
    logger.info(f"Pillow thumbnail saved → {out_path}")


# ---------------------------------------------------------------------------
# Resize PNG to web thumbnail
# ---------------------------------------------------------------------------

def make_thumbnail(src_png: str, out_path: str,
                   width: int = THUMB_W, height: int = THUMB_H) -> None:
    from PIL import Image
    img = Image.open(src_png)
    img = img.convert("RGB")
    img.thumbnail((width, height), Image.LANCZOS)
    # Pad to exact dimensions with background fill
    bg_color = (250, 250, 250)
    canvas = Image.new("RGB", (width, height), bg_color)
    x_off = (width - img.width) // 2
    y_off = (height - img.height) // 2
    canvas.paste(img, (x_off, y_off))
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    canvas.save(out_path, "PNG", optimize=True)
    logger.info(f"Thumbnail saved → {out_path}")


# ---------------------------------------------------------------------------
# Per-template preview pipeline
# ---------------------------------------------------------------------------

async def process_template(db, template, theme_obj, force: bool) -> None:
    from seeds.seed_runner import generate_ai_slides, build_premium_pptx

    name = template.name
    slug = name.lower().replace(" ", "_")

    pptx_path = PREVIEWS_DIR / f"{slug}_preview.pptx"
    image_path = IMAGES_DIR / f"{slug}_slide1.png"
    thumb_path = THUMBNAILS_DIR / f"{slug}_thumb.png"

    IMAGES_DIR.mkdir(parents=True, exist_ok=True)
    THUMBNAILS_DIR.mkdir(parents=True, exist_ok=True)

    theme_data = {"colors": theme_obj.colors, "fonts": theme_obj.fonts}

    # Step 1: Generate PPTX if missing or forced
    if not pptx_path.exists() or force:
        meta = template.metadata_json or {}
        slide_count = meta.get("total_slides", 8)
        logger.info(f"[{name}] Generating AI slides ({slide_count} slides)...")
        try:
            ai_slides = await generate_ai_slides(
                name,
                template.category or "",
                template.description or "",
                theme_data,
                slide_count,
            )
            build_premium_pptx(ai_slides, theme_data, str(pptx_path), name)
        except Exception as exc:
            logger.info(f"[{name}] PPTX generation failed: {exc}")
            # Fallback: write a minimal single-slide PPTX so COM export still works
            _write_fallback_pptx(name, theme_data, str(pptx_path))
    else:
        logger.info(f"[{name}] Reusing existing PPTX: {pptx_path}")

    # Step 2: Export slide 1 → PNG via PowerPoint COM
    exported_png = False
    if pptx_path.exists() and (not image_path.exists() or force):
        exported_png = pptx_slide_to_png(str(pptx_path), str(image_path), slide_index=0)

    # Step 3: Build web thumbnail
    if exported_png and image_path.exists():
        make_thumbnail(str(image_path), str(thumb_path))
    elif not thumb_path.exists() or force:
        # Pillow fallback — draw branded card
        logger.info(f"[{name}] Using Pillow fallback thumbnail")
        draw_thumbnail_from_theme(
            name,
            template.category or "Presentation",
            template.description or "",
            theme_obj.colors,
            str(thumb_path),
        )

    # Step 4: Persist paths to DB
    changed = False
    if pptx_path.exists() and str(template.preview_pptx_path) != str(pptx_path):
        template.preview_pptx_path = str(pptx_path)
        changed = True
    thumb_url = f"/previews/thumbnails/{slug}_thumb.png"
    if template.thumbnail_url != thumb_url:
        template.thumbnail_url = thumb_url
        changed = True
    if changed:
        await db.commit()
        logger.info(f"[{name}] DB updated — pptx_path + thumbnail_url")


def _write_fallback_pptx(name: str, theme_data: dict, out_path: str) -> None:
    """Write a minimal single-slide PPTX so COM still has something to export."""
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    def _rgb(h):
        h = h.lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    prs = PPTXPresentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb(theme_data["colors"].get("primary", "#0A2342"))

    txb = slide.shapes.add_textbox(Emu(400000), Emu(2000000), Emu(8000000), Emu(1000000))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = _rgb("#FFFFFF")

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

async def main(force: bool = False, only_template: str | None = None) -> None:
    import app.core.database_models  # noqa: F401
    await init_db()

    from sqlalchemy import select
    from app.models.template import Template
    from app.models.theme import Theme

    async with _db_module._session_factory() as db:
        query = select(Template).where(Template.is_active == True)  # noqa: E712
        templates = (await db.execute(query)).scalars().all()
        logger.info(f"Found {len(templates)} active templates")

        for template in templates:
            if only_template and template.name.lower() != only_template.lower():
                continue

            theme_obj = (
                await db.execute(select(Theme).where(Theme.id == template.theme_id))
            ).scalar_one_or_none()
            if not theme_obj:
                logger.info(f"[{template.name}] Theme not found, skipping")
                continue

            logger.info(f"\n{'='*60}")
            logger.info(f"Processing: {template.name}  ({template.category})")
            try:
                await process_template(db, template, theme_obj, force=force)
            except Exception as exc:
                logger.info(f"[{template.name}] ERROR: {exc}")

            logger.info("Waiting 3s before next template...")
            await asyncio.sleep(3)

    await close_db()
    logger.info("\nAll previews and thumbnails done.")
    logger.info(f"  PPTX files  → {PREVIEWS_DIR}/")
    logger.info(f"  Full images → {IMAGES_DIR}/")
    logger.info(f"  Thumbnails  → {THUMBNAILS_DIR}/")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate PPTX previews + PNG thumbnails")
    parser.add_argument("--force", action="store_true", help="Regenerate even if files exist")
    parser.add_argument("--template", type=str, default=None, help="Process only this template name")
    args = parser.parse_args()
    asyncio.run(main(force=args.force, only_template=args.template))
