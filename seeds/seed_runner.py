"""
Seed runner for WACDeckStudio.

Usage:
    python seeds/seed_runner.py [--force-rebuild]
"""
from __future__ import annotations

import argparse
import asyncio
import json
import sys
from pathlib import Path

# Add backend root to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.config import settings
from app.core.database import close_db, init_db
import app.core.database as _db_module
from app.utils.logger import get_logger

logger = get_logger("seed_runner")

SEEDS_DIR = Path(__file__).parent
THEMES_DIR = SEEDS_DIR / "themes"
TEMPLATES_DIR = SEEDS_DIR / "templates"
PREVIEWS_DIR = SEEDS_DIR / "previews"


# ---------------------------------------------------------------------------
# PPTX preview generation -- fully AI-generated, no placeholders
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))



async def generate_ai_slides(template_name: str, category: str, description: str,
                              theme_data: dict, slide_count: int) -> list[dict]:
    """Ask Gemini to generate a complete, professional slide deck JSON."""
    from app.ai import gemini_client

    colors = theme_data["colors"]
    prompt = f"""You are a world-class presentation designer creating a premium, professional slide deck.

Template: "{template_name}" ({category})
Description: {description}
Company: Nexus Solutions (a B2B SaaS company)
Theme colors: primary={colors.get('primary','#2563eb')}, accent={colors.get('accent','#f59e0b')}, background={colors.get('background','#ffffff')}

Generate exactly {slide_count} slides. Each slide must have REAL, executive-level content -- no placeholders.

Return a JSON array of slides:
[
  {{
    "order": 1,
    "type": "title|agenda|content|two_column|quote|stats|team|timeline|closing",
    "title": "Slide heading (max 8 words)",
    "subtitle": "Supporting line (max 15 words, optional)",
    "body": ["bullet point 1 (10-20 words)", "bullet point 2", "bullet point 3"],
    "stats": [{{"value": "47%", "label": "Cost Reduction"}}, ...],
    "quote": "A powerful quote here (20-40 words)",
    "attribution": "-- Name, Title",
    "image_keywords": "3-5 keywords for a relevant professional photo",
    "use_dark_bg": true or false
  }}
]

Rules:
- Title slide: compelling headline + tagline, use_dark_bg: true
- Agenda: list 5-6 section names as body bullets
- Content slides: 3-5 impactful bullet points with real data/facts
- Stats slides: 4 boxes with bold metrics (revenue, growth %, users, etc.)
- Quote slide: memorable executive quote
- Closing: strong CTA, use_dark_bg: true
- image_keywords must be specific and visual (e.g. "business team meeting modern office")

Respond with a valid JSON array ONLY. No markdown, no explanation."""

    return await gemini_client.generate_json(prompt, retries=2)


def build_premium_pptx(slides: list[dict], theme_data: dict, output_path: str, template_name: str) -> None:
    """Build a professional PPTX from AI-generated slide JSON using shapes, colors and text only."""
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    W_EMU = 9144000   # 16:9 width
    H_EMU = 5143500   # 16:9 height

    colors = theme_data["colors"]
    fonts  = theme_data["fonts"]

    primary_hex  = colors.get("primary",    "#2563eb")
    accent_hex   = colors.get("accent",     "#f59e0b")
    bg_hex       = colors.get("background", "#ffffff")
    text_hex     = colors.get("text",       "#1e293b")
    dark_bg_hex  = "#0f172a"

    heading_font   = fonts.get("heading", {}).get("family", "Calibri")
    body_font_name = fonts.get("body",    {}).get("family", "Calibri")

    prs = PPTXPresentation()
    prs.slide_width  = Emu(W_EMU)
    prs.slide_height = Emu(H_EMU)
    blank_layout = prs.slide_layouts[6]

    def _rgb(h: str) -> RGBColor:
        try:
            return hex_to_rgb(h)
        except Exception:
            return RGBColor(0x1e, 0x29, 0x3b)

    def _set_bg(slide, hex_color: str):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(hex_color)

    def _rect(slide, x, y, w, h, fill_hex: str):
        shape = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_hex)
        shape.line.fill.background()

    def _text(slide, text: str, x, y, w, h,
              font=None, size=18, bold=False, color="#ffffff",
              align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name  = font or body_font_name
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)

    for sd in slides:
        slide     = prs.slides.add_slide(blank_layout)
        use_dark  = sd.get("use_dark_bg", False)
        stype     = sd.get("type", "content")
        title_txt = sd.get("title", "")
        sub_txt   = sd.get("subtitle", "")
        body      = sd.get("body", [])
        stats     = sd.get("stats", [])
        quote_txt = sd.get("quote", "")
        attr_txt  = sd.get("attribution", "")

        bg_color   = dark_bg_hex if use_dark else bg_hex
        body_color = "#e2e8f0"   if use_dark else text_hex
        head_color = "#ffffff"   if use_dark else primary_hex

        _set_bg(slide, bg_color)

        # ?? Title ?????????????????????????????????????????????????????????
        if stype == "title":
            # Full left-side dark panel
            _rect(slide, 0, 0, W_EMU, H_EMU, dark_bg_hex)
            # Accent bar
            _rect(slide, 0, H_EMU // 2 - 25000, 500000, 50000, accent_hex)
            # Decorative right-side colour block
            _rect(slide, W_EMU - 1200000, 0, 1200000, H_EMU, primary_hex)
            _text(slide, title_txt,
                  300000, H_EMU // 2 - 750000, W_EMU - 1800000, 600000,
                  heading_font, 48, True, "#ffffff", PP_ALIGN.LEFT)
            if sub_txt:
                _text(slide, sub_txt,
                      300000, H_EMU // 2 - 100000, W_EMU - 1800000, 280000,
                      body_font_name, 24, False, "#94a3b8", PP_ALIGN.LEFT)
            _text(slide, template_name,
                  300000, H_EMU - 320000, 900000, 160000,
                  body_font_name, 13, False, "#64748b", PP_ALIGN.LEFT)

        # ?? Stats ?????????????????????????????????????????????????????????
        elif stype == "stats":
            _rect(slide, 0, 0, W_EMU, H_EMU, dark_bg_hex)
            _text(slide, title_txt,
                  300000, 160000, W_EMU - 600000, 360000,
                  heading_font, 36, True, "#ffffff", PP_ALIGN.LEFT)
            _rect(slide, 300000, 540000, 400000, 10000, accent_hex)
            stat_items = (stats or [])[:4]
            count = max(len(stat_items), 1)
            box_w = (W_EMU - 600000) // count
            for i, st in enumerate(stat_items):
                bx = 300000 + i * box_w
                _rect(slide, bx + 40000, 600000, box_w - 80000, 1300000, primary_hex)
                _text(slide, st.get("value", "--"),
                      bx + 60000, 680000, box_w - 120000, 550000,
                      heading_font, 54, True, "#ffffff", PP_ALIGN.CENTER)
                _text(slide, st.get("label", ""),
                      bx + 60000, 1260000, box_w - 120000, 280000,
                      body_font_name, 16, False, accent_hex, PP_ALIGN.CENTER)

        # ?? Quote ?????????????????????????????????????????????????????????
        elif stype == "quote":
            _rect(slide, 0, 0, W_EMU, H_EMU, dark_bg_hex)
            _rect(slide, 0, 0, 18000, H_EMU, accent_hex)          # left accent stripe
            _text(slide, '"',
                  200000, 80000, 700000, 700000,
                  heading_font, 120, True, accent_hex, PP_ALIGN.LEFT)
            _text(slide, quote_txt,
                  300000, 580000, W_EMU - 600000, 1500000,
                  heading_font, 28, False, "#f1f5f9", PP_ALIGN.CENTER, italic=True)
            if attr_txt:
                _text(slide, attr_txt,
                      300000, H_EMU - 480000, W_EMU - 600000, 220000,
                      body_font_name, 18, True, accent_hex, PP_ALIGN.CENTER)

        # ?? Closing ???????????????????????????????????????????????????????
        elif stype == "closing":
            _rect(slide, 0, 0, W_EMU, H_EMU, dark_bg_hex)
            _rect(slide, W_EMU // 2 - 600000, H_EMU // 2 - 15000, 1200000, 12000, accent_hex)
            _text(slide, title_txt,
                  300000, H_EMU // 2 - 720000, W_EMU - 600000, 540000,
                  heading_font, 48, True, "#ffffff", PP_ALIGN.CENTER)
            if sub_txt:
                _text(slide, sub_txt,
                      300000, H_EMU // 2 - 140000, W_EMU - 600000, 280000,
                      body_font_name, 22, False, "#94a3b8", PP_ALIGN.CENTER)

        # ?? Agenda / two_column / timeline ? treat as content ?????????????
        # ?? Content (default) ?????????????????????????????????????????????
        else:
            # Left accent bar
            _rect(slide, 300000, 260000, 14000, H_EMU - 520000, accent_hex)
            _text(slide, title_txt,
                  360000, 260000, W_EMU - 720000, 460000,
                  heading_font, 36, True, head_color, PP_ALIGN.LEFT)
            # Underline rule
            _rect(slide, 360000, 740000, 700000, 8000, primary_hex)

            y = 810000
            for bullet in body:
                clean = bullet.lstrip("?- ").strip()
                if not clean:
                    continue
                _text(slide, f"?  {clean}",
                      380000, y, W_EMU - 760000, 340000,
                      body_font_name, 18, False, body_color, PP_ALIGN.LEFT)
                y += 360000
                if y > H_EMU - 380000:
                    break

            if sub_txt and not body:
                _text(slide, sub_txt,
                      380000, 840000, W_EMU - 760000, 300000,
                      body_font_name, 22, False, "#64748b", PP_ALIGN.LEFT)

        # Slide number (bottom-right)
        _text(slide, str(sd.get("order", "")),
              W_EMU - 420000, H_EMU - 260000, 280000, 180000,
              body_font_name, 11, False, "#94a3b8", PP_ALIGN.RIGHT)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    logger.info(f"PPTX saved ? {output_path}")


# ---------------------------------------------------------------------------
# Seed themes
# ---------------------------------------------------------------------------

async def seed_themes(db, force_rebuild: bool = False) -> dict[str, str]:
    """Upsert all themes. Returns name->id mapping."""
    from sqlalchemy import select
    from app.models.theme import Theme

    theme_id_map: dict[str, str] = {}

    for theme_file in THEMES_DIR.glob("*.json"):
        data = json.loads(theme_file.read_text(encoding="utf-8"))
        name = data["name"]

        existing = (await db.execute(select(Theme).where(Theme.name == name))).scalar_one_or_none()

        if existing and not force_rebuild:
            logger.info(f"Theme '{name}' already exists, skipping")
            theme_id_map[name] = str(existing.id)
            continue

        if existing and force_rebuild:
            existing.colors = data["colors"]
            existing.fonts = data["fonts"]
            await db.commit()
            await db.refresh(existing)
            theme_id_map[name] = str(existing.id)
            logger.info(f"Theme '{name}' updated")
        else:
            theme = Theme(
                name=name,
                colors=data["colors"],
                fonts=data["fonts"],
            )
            db.add(theme)
            await db.commit()
            await db.refresh(theme)
            theme_id_map[name] = str(theme.id)
            logger.info(f"Theme '{name}' created")

    return theme_id_map


# ---------------------------------------------------------------------------
# Seed templates + generate PPTX previews
# ---------------------------------------------------------------------------

async def fill_placeholders(template_data: dict) -> dict:
    """Replace [PLACEHOLDER: ...] tokens with realistic sample content (no API calls needed)."""
    import copy
    import re

    COMPANY = "Nexus Solutions"
    YEAR = "2025"

    # Keyword-based replacement table -- ordered from most specific to least
    RULES: list[tuple[str, str]] = [
        # Company / branding
        ("company name", COMPANY),
        ("company", COMPANY),
        ("brand name", COMPANY),
        ("product name", "NexusCore Platform"),
        ("product", "NexusCore Platform"),
        ("app name", "NexusCore"),
        # Time
        ("year", YEAR),
        ("quarter", "Q3 2025"),
        ("date", "October 2025"),
        ("timeline", "Q4 2025 -- Q2 2026"),
        ("deadline", "December 31, 2025"),
        ("launch date", "November 15, 2025"),
        ("duration", "6 Months"),
        # People
        ("ceo", "Sarah Mitchell"),
        ("founder", "Sarah Mitchell"),
        ("presenter", "Sarah Mitchell"),
        ("team lead", "James Carter"),
        ("name", "Alex Rivera"),
        ("role", "Senior Product Manager"),
        ("focus area", "Growth & Retention"),
        ("job", "Marketing Director"),
        ("age", "32"),
        ("pain", "Too much time on manual workflows"),
        # Titles / headings
        ("subtitle or theme", f"Scaling Smarter in {YEAR}"),
        ("subtitle", f"Building the Future Together"),
        ("tagline", "Efficiency. Innovation. Growth."),
        ("title", f"{COMPANY} -- {YEAR} Strategy"),
        ("section heading", "Strategic Priorities"),
        ("heading", "Our Path Forward"),
        # Goals / targets
        ("grow mqls", "Grow MQLs by 40%"),
        ("increase brand awareness", "Increase brand awareness by 35%"),
        ("reduce cac", "Reduce CAC by 25%"),
        ("goal 1", "Expand enterprise segment by 50%"),
        ("goal 2", "Launch 3 new product tiers"),
        ("goal 3", "Achieve 95% customer retention"),
        ("goal 4", "Enter 2 new international markets"),
        ("annual target", "200% ARR Growth"),
        ("target", "$4.2M ARR"),
        ("objective", "Achieve market leadership in SMB segment"),
        # Financial
        ("total budget", "$1,200,000"),
        ("budget", "$850,000"),
        ("revenue", "$12.4M"),
        ("arr", "$8.7M ARR"),
        ("mrr", "$725K MRR"),
        ("cost", "$320,000"),
        ("roi", "320% ROI"),
        ("funding", "$15M Series A"),
        ("valuation", "$85M"),
        # Metrics / KPIs
        ("kpi 1", "Customer Acquisition Cost\nTarget: $120"),
        ("kpi 2", "Monthly Recurring Revenue\nTarget: $900K"),
        ("kpi 3", "Net Promoter Score\nTarget: 72+"),
        ("kpi 4", "Churn Rate\nTarget: < 2%"),
        ("kpi", "Monthly Active Users -- 45,000"),
        ("metric", "Customer Lifetime Value: $4,800"),
        ("conversion rate", "8.4% Conversion Rate"),
        ("growth rate", "127% YoY Growth"),
        ("market size", "$42B TAM"),
        ("percentage", "68%"),
        ("number", "12,500"),
        # Marketing
        ("channel 1", "Paid Search (Google Ads)"),
        ("channel 2", "Content Marketing & SEO"),
        ("channel 3", "LinkedIn Outreach"),
        ("channel 4", "Email Automation"),
        ("channel mix", "Omnichannel Demand Generation"),
        ("channel", "Digital & Social Media"),
        ("campaign name", "Nexus Growth Sprint"),
        ("campaign", "Q3 Pipeline Acceleration"),
        ("content pillar 1", "Thought Leadership & Education"),
        ("content pillar 2", "Product Stories & Case Studies"),
        ("content pillar 3", "Community & Partner Spotlights"),
        ("content pillar", "Data-Driven Insights"),
        ("persona 1", "Persona 1: The Ops Leader\nAge: 38\nJob: VP Operations\nPain: Manual processes slowing growth"),
        ("persona 2", "Persona 2: The Growth Marketer\nAge: 29\nJob: Digital Marketing Manager\nPain: Fragmented data across tools"),
        ("persona 3", "Persona 3: The CFO\nAge: 45\nJob: Chief Financial Officer\nPain: Lack of real-time financial visibility"),
        ("audience", "Enterprise B2B -- 500+ employees"),
        ("target audience", "Mid-Market SaaS Companies"),
        # SWOT
        ("strengths", "Strong product-market fit\nExperienced founding team\n92% customer satisfaction score"),
        ("weaknesses", "Limited brand awareness\nSmall sales team relative to opportunity\nNarrow geographic presence"),
        ("opportunities", "$42B untapped market\nCompetitor product gaps\nAI/automation tailwinds accelerating"),
        ("threats", "Well-funded incumbents\nRegulatory changes in data privacy\nTalent market competition"),
        ("where we stand", "Positioned for Hypergrowth"),
        # Budget allocation
        ("digital ads", "Paid Digital -- 35%"),
        ("content creation", "Content & Creative -- 20%"),
        ("events", "Events & Conferences -- 15%"),
        ("tools & tech", "Marketing Tech Stack -- 30%"),
        ("alloc", "Paid Media -- $420,000 (35%)"),
        # Team
        ("team", f"{COMPANY} Leadership Team"),
        ("responsibilities", "Own the full go-to-market strategy"),
        ("focus areas", "Pipeline Generation, ABM, Brand"),
        # Reporting
        ("reporting cadence", "Weekly dashboards ? Monthly business reviews ? Quarterly OKR check-ins"),
        ("reporting", "Bi-weekly review cadence with leadership"),
        # CTA / closing
        ("let's grow together", f"Let's Build the Future Together"),
        ("contact", f"hello@nexussolutions.io  ?  +1 (415) 867-5309"),
        ("email", "hello@nexussolutions.io"),
        ("website", "www.nexussolutions.io"),
        ("next steps", "Schedule discovery call ? Review proposal ? Sign agreement"),
        ("call to action", "Book Your Free Demo Today"),
        # Problem / solution
        ("problem", "Teams waste 40% of their workweek on manual, repetitive tasks"),
        ("solution", "NexusCore automates workflows end-to-end, cutting ops costs by 60%"),
        ("value proposition", "10x faster workflows. Zero manual errors."),
        ("pain point", "Disconnected tools create data silos and slow decision-making"),
        ("differentiator", "Only platform combining AI automation with real-time analytics"),
        ("competitive advantage", "Proprietary ML engine trained on 50M+ enterprise workflows"),
        # Roadmap / project
        ("phase 1", "Discovery & Architecture -- 6 weeks"),
        ("phase 2", "Development & Integration -- 12 weeks"),
        ("phase 3", "Testing & Launch -- 4 weeks"),
        ("milestone", "MVP Launch -- November 15, 2025"),
        ("deliverable", "Fully integrated enterprise dashboard"),
        ("risk", "Integration complexity with legacy systems"),
        ("mitigation", "Dedicated integration team with 3-week buffer"),
        # Generic fallbacks
        ("description", f"A next-generation platform for enterprise automation"),
        ("summary", "Nexus Solutions delivers measurable results through intelligent automation"),
        ("insight", "83% of enterprise teams cite manual processes as their #1 growth barrier"),
        ("stat", "47% cost reduction within first 90 days"),
        ("quote", '"The platform transformed how our entire operations team works." -- COO, Fortune 500 Co.'),
        ("benefit", "Reduces operational overhead by 60% in under 3 months"),
        ("feature", "Real-time analytics with AI-powered recommendations"),
    ]

    slides_json = json.dumps(template_data.get("slides", []))
    tokens = list(dict.fromkeys(re.findall(r'\[PLACEHOLDER[^\]]*\]', slides_json)))
    if not tokens:
        return template_data

    def _match(token: str) -> str:
        desc = re.sub(r'\[PLACEHOLDER[:\s]*', '', token).rstrip(']').strip().lower()
        for keyword, replacement in RULES:
            if keyword in desc:
                return replacement
        # last-resort: title-case the description
        return desc.replace(' ', ' ').title() if desc else COMPANY

    replacements = {tok: _match(tok) for tok in tokens}

    filled = copy.deepcopy(template_data)
    filled_str = json.dumps(filled)
    for token, replacement in replacements.items():
        filled_str = filled_str.replace(token, replacement.replace('\n', '\\n'))
    try:
        filled = json.loads(filled_str)
        tname = template_data.get('name', '?')
        logger.info(f"Placeholders filled for '{tname}' ({len(replacements)} tokens)")
    except json.JSONDecodeError:
        return template_data

    return filled


async def seed_templates(db, theme_id_map: dict[str, str], force_rebuild: bool = False, skip_previews: bool = False) -> None:
    """Upsert all templates and generate PPTX previews."""
    from sqlalchemy import select
    from app.models.template import Template
    from app.models.theme import Theme

    for tmpl_file in TEMPLATES_DIR.glob("*.json"):
        raw = tmpl_file.read_text(encoding="utf-8")
        data = json.loads(raw)
        name = data["name"]
        theme_name = data.pop("theme_name", None)

        if theme_name is None:
            logger.info(f"Template '{name}' missing theme_name, skipping")
            continue

        theme_id_str = theme_id_map.get(theme_name)
        if not theme_id_str:
            logger.info(f"Theme '{theme_name}' not found for template '{name}', skipping")
            continue

        metadata = data.get("metadata", {})
        actual_slide_count = len(data.get("slides", []))
        if metadata.get("total_slides", 0) != actual_slide_count:
            logger.info(
                f"WARNING: Template '{name}' metadata.total_slides={metadata.get('total_slides')} "
                f"but has {actual_slide_count} slides -- using actual count"
            )
            metadata["total_slides"] = actual_slide_count

        existing = (await db.execute(select(Template).where(Template.name == name))).scalar_one_or_none()

        # Generate preview PPTX
        preview_filename = name.lower().replace(" ", "_") + "_preview.pptx"
        preview_path = str(PREVIEWS_DIR / preview_filename)

        # Load theme data for preview generation
        theme_obj = (
            await db.execute(select(Theme).where(Theme.id == theme_id_str))
        ).scalar_one_or_none()
        theme_data_for_preview = {"colors": theme_obj.colors, "fonts": theme_obj.fonts} if theme_obj else None

        if skip_previews:
            preview_path = None
        elif theme_data_for_preview:
            try:
                slide_count = metadata.get("total_slides", 8)
                logger.info(f"Generating AI slides for '{name}' ({slide_count} slides)...")
                ai_slides = await generate_ai_slides(
                    name, data.get("category", ""), data.get("description", ""),
                    theme_data_for_preview, slide_count,
                )
                build_premium_pptx(ai_slides, theme_data_for_preview, preview_path, name)
            except Exception as exc:
                logger.info(f"AI preview generation failed for '{name}': {exc}")
                preview_path = None
        else:
            preview_path = None

        if existing and not force_rebuild:
            # Update preview path if we just generated it
            if preview_path and existing.preview_pptx_path != preview_path:
                existing.preview_pptx_path = preview_path
                await db.commit()
            logger.info(f"Template '{name}' already exists, skipping data update")
            continue

        if existing and force_rebuild:
            existing.description = data.get("description", "")
            existing.category = data.get("category", "")
            existing.tags = data.get("tags", [])
            existing.theme_id = theme_id_str
            existing.metadata_json = metadata
            existing.slides = data.get("slides", [])
            existing.is_active = data.get("is_active", True)
            existing.thumbnail_url = data.get("thumbnail_url", "")
            existing.preview_pptx_path = preview_path
            await db.commit()
            logger.info(f"Template '{name}' updated")
        else:
            template = Template(
                name=name,
                description=data.get("description", ""),
                category=data.get("category", ""),
                tags=data.get("tags", []),
                thumbnail_url=data.get("thumbnail_url", ""),
                theme_id=theme_id_str,
                is_active=data.get("is_active", True),
                metadata_json=metadata,
                slides=data.get("slides", []),
                preview_pptx_path=preview_path,
            )
            db.add(template)
            await db.commit()
            logger.info(f"Template '{name}' created")


# ---------------------------------------------------------------------------
# Seed sample user
# ---------------------------------------------------------------------------

async def seed_user(db) -> None:
    from sqlalchemy import select
    from app.models.user import User
    from app.core.security import hash_password

    email = "demo@wacdeckstudio.com"
    existing = (await db.execute(select(User).where(User.email == email))).scalar_one_or_none()
    if existing:
        logger.info(f"User '{email}' already exists, skipping")
        return

    user = User(
        email=email,
        hashed_password=hash_password("demo1234"),
        full_name="Demo User",
        role="user",
        is_active=True,
    )
    db.add(user)
    await db.commit()
    logger.info(f"Sample user created -- email: {email}  password: demo1234")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

async def main(force_rebuild: bool = False) -> None:
    logger.info("Initializing database connection...")
    import app.core.database_models  # noqa: F401
    await init_db()

    async with _db_module._session_factory() as db:
        logger.info("Seeding sample user...")
        await seed_user(db)

        logger.info("Seeding themes...")
        theme_id_map = await seed_themes(db, force_rebuild=force_rebuild)

        logger.info("Seeding templates (with PPTX preview generation)...")
        await seed_templates(db, theme_id_map, force_rebuild=force_rebuild)

    await close_db()
    logger.info("Seed complete!")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Seed WACDeckStudio database")
    parser.add_argument("--force-rebuild", action="store_true", help="Force update existing records")
    args = parser.parse_args()
    asyncio.run(main(force_rebuild=args.force_rebuild))
