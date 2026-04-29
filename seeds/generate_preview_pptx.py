"""
Generate Gamma-quality preview PPTX for each template seed.

Usage:
    cd backend
    python seeds/generate_preview_pptx.py
    python seeds/generate_preview_pptx.py --force
    python seeds/generate_preview_pptx.py --template business_pitch
"""
from __future__ import annotations

import argparse
import asyncio
import io
import json
import re as _re
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

TEMPLATES_DIR = Path(__file__).parent / "templates"
THEMES_DIR    = Path(__file__).parent / "themes"
PREVIEWS_DIR  = Path(__file__).parent / "previews"
LOGO_PATH     = Path(__file__).resolve().parent.parent / "logo" / "logo1.svg"

DESIGN_W    = 1920
DESIGN_H    = 1080
SLIDE_W_IN  = 13.33
SLIDE_H_IN  = 7.5
EMU_PER_INCH = 914400

NEUTRAL_BACKGROUNDS = {"#F5F5F5", "#F0F0F0", "#FAFAFA", "#E8EAF6"}

SEED_DARK_MAP = {
    "#1A1A2E": "primary",
    "#16213E": "secondary",
    "#0F3460": "accent",
}

DARK_TEXT_COLOR = "#1A1A2E"
CARD_PADDING_PX = 40

# ─── Fallback demo content (used when Gemini quota is exhausted) ──────────────

FALLBACK: dict[str, dict[str, str]] = {
    "business_pitch": {
        "s1-title":      "Nexus AI",
        "s1-subtitle":   "Eliminating manual ops with intelligent automation",
        "s1-date":       "April 2026",
        "s2-heading":    "Operations teams waste 12 hours a week on manual work",
        "s2-desc":       "Finance, HR, and RevOps teams manually reconcile data across 6+ tools every day. Errors cost companies an average of $320K per year in rework and lost deals.",
        "s2-stat1":      "12 hrs\nlost / week",
        "s2-stat2":      "$320K\ncost per year",
        "s3-heading":    "One platform. Zero manual work.",
        "s3-feat1":      "AI Workflows\nAutomate any process in minutes",
        "s3-feat2":      "Smart Integrations\n200+ native connectors",
        "s3-feat3":      "Live Analytics\nReal-time ops visibility",
        "s4-heading":    "$18B market growing 34% YoY",
        "s4-tam":        "TAM\n$18B",
        "s4-sam":        "SAM\n$4.2B",
        "s4-som":        "SOM\n$420M",
        "s4-growth":     "The workflow automation market is projected to reach $26B by 2028, driven by AI adoption in mid-market enterprises.",
        "s5-heading":    "Up and running in 15 minutes",
        "s5-step1":      "1. Connect\nyour tools",
        "s5-step2":      "2. Define\nyour workflow",
        "s5-step3":      "3. AI builds\nthe automation",
        "s5-step4":      "4. Go live\n& monitor",
        "s6-heading":    "2,400 teams. $12M ARR. 127% NRR.",
        "s6-metric1":    "2,400\nActive Users",
        "s6-metric2":    "$1.0M\nMRR",
        "s6-metric3":    "22%\nMoM Growth",
        "s6-metric4":    "47\nEnterprise Partners",
        "s7-heading":    "SaaS + professional services",
        "s7-stream1":    "Platform Subscription\nStarter $299/mo · Growth $899/mo · Enterprise custom",
        "s7-stream2":    "Managed Automation\nDedicated CS + build-outs from $2,500/mo",
        "s7-unit-econ":  "LTV $28K · CAC $3,200 · Payback 4 months · NRR 127%",
        "s8-heading":    "We win on speed, depth, and AI",
        "s8-differentiator": "Unlike Zapier (simple triggers) or UiPath (RPA, IT-heavy), Nexus AI provides business-user-friendly AI workflows with enterprise-grade security and instant ROI.",
        "s8-matrix":     "Feature comparison vs Zapier, UiPath, Make — Nexus AI leads on AI-native, time-to-value, and SMB pricing",
        "s9-heading":    "Meet the founders",
        "s9-member1":    "Priya Shah\nCEO\nFormer VP Ops @ Stripe, MIT MBA",
        "s9-member2":    "Lucas Reeves\nCTO\nEx-Google Brain, 12 patents in ML",
        "s9-member3":    "Amara Osei\nCRO\nBuilt 0→$50M ARR @ Ramp",
        "s10-heading":   "18-month roadmap",
        "s10-q1":        "Q1\nLaunch enterprise tier",
        "s10-q2":        "Q2\nEU data residency + SOC 2",
        "s10-q3":        "Q3\nAI co-pilot GA",
        "s10-q4":        "Q4+\nIPO readiness",
        "s11-heading":   "Series A — $8M",
        "s11-ask":       "$8M Series A",
        "s11-use1":      "Product R&D — 45%",
        "s11-use2":      "Sales & Marketing — 35%",
        "s11-use3":      "Ops & G&A — 20%",
        "s12-cta":       "Let's Build the Future Together",
        "s12-contact":   "hello@nexusai.io  ·  www.nexusai.io",
    },
    "startup_deck": {
        "s1-title":    "Velo",
        "s1-tagline":  "Instant global payments. Zero friction.",
        "s2-heading":  "Our Vision",
        "s2-vision":   "A world where moving money is as fast and free as sending a text message.",
        "s2-why":      "Real-time rails are live in 40+ countries. The infrastructure to connect them is missing. That's Velo.",
        "s3-heading":  "The Payment Problem",
        "s3-pain1":    "Cross-border transfers take 3–5 days",
        "s3-pain2":    "Hidden FX fees eat 3–7% of every transaction",
        "s3-pain3":    "Reconciliation is entirely manual for 80% of SMBs",
        "s4-heading":  "Velo Pay",
        "s4-desc":     "A single API connecting 40 real-time payment networks. Send money anywhere in under 8 seconds at cost-plus pricing.",
        "s4-benefit1": "Sub-8s settlement globally",
        "s4-benefit2": "0.3% flat FX — no hidden fees",
        "s4-benefit3": "Auto-reconciliation via AI",
        "s5-heading":  "Traction speaks for itself",
        "s5-num1":     "1,200+\nCustomers",
        "s5-num2":     "$480K\nARR",
        "s5-num3":     "31%\nGrowth MoM",
        "s7-heading":  "World-class team",
        "s7-m1":       "Dan Carter\nCEO — Ex-Wise, Forbes 30U30",
        "s7-m2":       "Yuna Kim\nCTO — Ex-Stripe Engineering Lead",
        "s7-m3":       "Raj Patel\nCMO — Built growth @ Revolut APAC",
        "s8-heading":  "A $220B opportunity",
        "s8-tam":      "TAM\n$220B",
        "s8-sam":      "SAM\n$38B",
        "s8-som":      "SOM\n$950M",
        "s9-heading":  "Join our journey",
        "s9-amount":   "$3.5M Seed Round",
        "s9-use1":     "55% Product & Engineering",
        "s9-use2":     "30% Growth & Marketing",
        "s9-use3":     "15% Team Expansion",
        "s10-cta":     "Let's Change How the World Moves Money",
        "s10-contact": "dan@velopay.io  ·  www.velopay.io",
    },
    "product_launch": {
        "s1-product":  "Luma",
        "s1-tagline":  "AI project management that actually thinks ahead",
        "s2-heading":  "Project management is broken",
        "s2-p1":       "Teams spend 4 hours/week updating status in PM tools",
        "s2-p2":       "60% of projects miss deadlines due to poor visibility",
        "s2-p3":       "Context switching across 5+ tools kills deep work",
        "s2-p4":       "Reporting to stakeholders is a manual nightmare",
        "s3-heading":  "Luma Changes Everything",
        "s3-desc":     "Luma uses AI to auto-update tasks, predict blockers 48 hours in advance, and generate stakeholder reports in one click.",
        "s3-feat1":    "Auto-updates\nfrom Slack & email",
        "s3-feat2":    "Blocker prediction\n48h ahead",
        "s3-feat3":    "1-click reports\nfor any audience",
        "s4-heading":  "See It in Action",
        "s5-heading":  "Why teams love Luma",
        "s5-b1":       "Save 4 hrs/week — no more manual status updates",
        "s5-b2":       "Hit 94% of deadlines with AI-powered early warnings",
        "s5-b3":       "Every stakeholder always in sync, zero effort",
        "s5-b4":       "Connects to Jira, Linear, Notion, Slack in 2 minutes",
        "s6-heading":  "What beta users say",
        "s6-q1":       "\"Luma replaced our entire Monday.com workflow AND our weekly standup.\"\n— Sarah T., Head of Product @ Findr",
        "s6-q2":       "\"The blocker prediction alone saved our Q3 launch. Absolute game-changer.\"\n— Marcus L., CTO @ Cresta",
        "s7-heading":  "Simple, transparent pricing",
        "s7-plan1":    "Starter\n$12/user/mo\nUp to 10 users\nCore AI features",
        "s7-plan2":    "Growth (Popular)\n$29/user/mo\nUnlimited users\nAdvanced AI + reports",
        "s7-plan3":    "Enterprise\nCustom\nSSO + SLA\nDedicated CS",
        "s8-heading":  "What's coming next",
        "s8-now":      "NOW\nAI task updates · Blocker alerts · 1-click reports",
        "s8-next":     "NEXT\nCalendar AI · Resource forecasting · Custom dashboards",
        "s8-later":    "LATER\nAutonomous sprints · Portfolio AI · Revenue intelligence",
        "s9-heading":  "Common questions",
        "s9-q1":       "Q: Does Luma replace our PM tool?\nA: No — it supercharges it. Luma sits on top of Jira, Linear, or Asana.",
        "s9-q2":       "Q: How does AI get context?\nA: It reads your Slack, email threads, and task history — no manual input.",
        "s9-q3":       "Q: Is our data secure?\nA: SOC 2 Type II certified. Data never used for model training.",
        "s9-q4":       "Q: How long to set up?\nA: Most teams are live in under 20 minutes.",
        "s10-heading": "Try Luma Free for 14 Days",
        "s10-cta":     "www.lumapm.ai/start",
    },
    "quarterly_review": {
        "s1-quarter":        "Q1 2026 BUSINESS REVIEW",
        "s1-title":          "Apex Corp\nQuarterly Review",
        "s1-presenter":      "Presented by: James Okafor, CEO",
        "s2-heading":        "Q1 at a Glance",
        "s2-win1":           "WIN: Closed largest enterprise deal in company history — $2.4M ARR",
        "s2-win2":           "WIN: Launched v3.0 on schedule — 99.98% uptime",
        "s2-challenge":      "CHALLENGE: EMEA expansion delayed 6 weeks due to regulatory review",
        "s2-overall":        "Overall vs Target: 104%",
        "s3-heading":        "Revenue This Quarter",
        "s3-total":          "$4.8M\nTotal Revenue",
        "s3-target":         "104% of Target",
        "s3-growth":         "+31% QoQ Growth",
        "s3-breakdown":      "Enterprise 62% · Mid-market 28% · SMB 10% · APAC +44% fastest growing region",
        "s4-heading":        "Performance Dashboard",
        "s4-kpi1":           "Net Revenue\n$4.8M\nvs $4.6M target",
        "s4-kpi2":           "New ARR\n$1.2M\nvs $1.1M target",
        "s4-kpi3":           "NRR\n128%\nvs 120% target",
        "s4-kpi4":           "Churn\n1.2%\nvs 2.0% target",
        "s4-trend":          "All four core KPIs beat Q1 targets. NRR hit an all-time high of 128%, driven by expansion in the enterprise segment.",
        "s5-heading":        "Sales Highlights",
        "s5-new":            "New Deals Closed\n47",
        "s5-pipeline":       "Pipeline Value\n$18.4M",
        "s5-conv":           "Conversion Rate\n24%",
        "s5-insight":        "Top performer: APAC team closed 3 deals >$500K. Enterprise sales cycle shortened from 92 to 71 days.",
        "s6-heading":        "What We Shipped",
        "s6-shipped1":       "v3.0 platform launch — zero-downtime migration for 2,100 customers",
        "s6-shipped2":       "AI anomaly detection in analytics module (90-day beta)",
        "s6-shipped3":       "Mobile app redesign — 4.8★ App Store rating",
        "s6-shipped4":       "Salesforce bi-directional sync (most-requested feature, 18 months)",
        "s7-heading":        "Customer Health",
        "s7-nps":            "NPS Score\n68",
        "s7-churn":          "Churn Rate\n1.2%",
        "s7-csat":           "CSAT\n94%",
        "s7-insight":        "3 at-risk accounts proactively saved via CS intervention. Spotlight: Meridian Corp expanded 3× seat count after v3.0 launch.",
        "s8-heading":        "Team Update",
        "s8-headcount":      "Team Size\n184 FTE",
        "s8-hires":          "New Hires\n22",
        "s8-open":           "Open Roles\n14",
        "s8-culture":        "Engineering expanded in Lisbon (8 hires). 'All-hands roadshow' in 3 cities — 97% attendance. Glassdoor rating 4.6★.",
        "s9-heading":        "What We Learned",
        "s9-c1":             "Challenge: EMEA regulatory delay\nAction: Retained local compliance counsel, dedicated EMEA GTM lead",
        "s9-c2":             "Challenge: Support ticket volume +40% post v3.0 launch\nAction: Expanded support team by 6, launched AI-assist for Tier-1",
        "s9-c3":             "Challenge: SMB churn uptick in Jan\nAction: Launched SMB success playbook, reduced churn to 1.2% by March",
        "s10-heading":       "Q2 Priorities",
        "s10-p1":            "Priority 1\nClose EMEA regulatory approval · target: May 15",
        "s10-p2":            "Priority 2\nLaunch AI anomaly detection GA · target: June 1",
        "s10-p3":            "Priority 3\nHit $5.4M revenue · +12.5% QoQ",
        "s10-revenue-target":"Q2 Revenue Target\n$5.4M",
        "s11-heading":       "Thank You — Questions Welcome",
        "s11-contact":       "james.okafor@apexcorp.io",
    },
    "project_proposal": {
        "s1-title":       "Digital Customer Portal",
        "s1-date":        "Submitted by: Elena Markov, Director of Digital  ·  April 2026",
        "s2-heading":     "Project Overview",
        "s2-summary":     "Build a self-service customer portal to deflect 40% of support tickets, reduce onboarding time from 14 to 5 days, and improve NPS by 20 points. Estimated 14-month ROI of 380%.",
        "s2-budget":      "Budget: $240,000",
        "s2-timeline":    "Timeline: 6 months",
        "s2-roi":         "ROI: 380%",
        "s3-heading":     "The Problem Today",
        "s3-desc":        "Customers currently email or call support for tasks they could self-serve: resetting credentials, viewing invoices, checking order status, and updating billing. This creates 1,800+ tickets/month at a fully-loaded cost of $22 per ticket.",
        "s3-impact":      "Annual cost of inaction: $475,000 in support overhead, plus 23% of churned customers cite 'poor self-service' as primary reason.",
        "s3-urgency":     "High Urgency ⚠",
        "s4-heading":     "What We Will Achieve",
        "s4-obj1":        "1. Deflect 40% of support tickets via self-service (save $190K/year)",
        "s4-obj2":        "2. Reduce onboarding time from 14 to 5 days (improve activation 35%)",
        "s4-obj3":        "3. Increase NPS from 42 to 62 within 6 months of launch",
        "s4-obj4":        "4. Achieve full feature parity with competitor portals by Q3 2026",
        "s5-heading":     "How We'll Do It",
        "s5-phase1":      "Phase 1\nDiscovery & Planning\nUser research, requirements, vendor selection",
        "s5-phase2":      "Phase 2\nDesign & Development\nUX design, API build, integrations",
        "s5-phase3":      "Phase 3\nTesting & Review\nBeta with 200 customers, load testing",
        "s5-phase4":      "Phase 4\nLaunch & Evaluation\nPhased rollout, KPI review at 30/60/90 days",
        "s6-heading":     "6-Month Roadmap",
        "s6-milestone1":  "Month 1–2: Discovery complete, vendor selected, design signed off",
        "s6-milestone2":  "Month 3–4: Core portal live in staging, beta customer group onboarded",
        "s6-milestone3":  "Month 5–6: GA launch, support deflection tracking live",
        "s6-deadline":    "Target Completion\nOctober 2026",
        "s7-heading":     "Investment Required",
        "s7-total":       "Total Budget\n$240,000",
        "s7-line1":       "Design & UX — $40,000",
        "s7-line2":       "Engineering (internal + contract) — $120,000",
        "s7-line3":       "Third-party integrations & licenses — $45,000",
        "s7-line4":       "QA, launch & contingency — $35,000",
        "s8-heading":     "Who's Involved",
        "s8-m1":          "Elena Markov\nProject Sponsor\nBudget owner, executive escalation",
        "s8-m2":          "Carlos Diaz\nEngineering Lead\nArchitecture, build, API design",
        "s8-m3":          "Aisha Brooks\nProduct Manager\nRequirements, roadmap, user research",
        "s8-m4":          "Tom Nguyen\nCS Lead\nBeta coordination, support integration",
        "s9-heading":     "Risk Management",
        "s9-r1":          "Risk: API integration delays with billing system\nImpact: High\nMitigation: Parallel track with fallback manual workflow",
        "s9-r2":          "Risk: Low beta adoption\nImpact: Medium\nMitigation: Incentivised early-access program with CS team support",
        "s9-r3":          "Risk: Scope creep\nImpact: Medium\nMitigation: Change control board, weekly steering committee",
        "s10-heading":    "Request for Approval",
        "s10-contact":    "Contact: elena.markov@company.com  ·  Ext. 4421",
    },
    "marketing_plan": {
        "s1-year":        "2026 MARKETING PLAN",
        "s1-title":       "Bloom\nGrowth Strategy",
        "s1-subtitle":    "Sustainable living, exponential reach",
        "s2-heading":     "Where We Stand",
        "s2-s":           "STRENGTHS\nStrong brand loyalty · 4.9★ avg product rating · 140K social followers",
        "s2-w":           "WEAKNESSES\nLow brand awareness outside core demo · Thin paid media capability",
        "s2-o":           "OPPORTUNITIES\nESG investment surge · Creator economy partnerships · EU expansion",
        "s2-t":           "THREATS\nGreenwashing scrutiny · Rising Meta CPMs · Amazon private label",
        "s3-heading":     "Who We're Reaching",
        "s3-persona1":    "Persona 1: The Conscious Millennial\nAge: 28–38\nJob: Urban professional\nMotivation: Values alignment + convenience",
        "s3-persona2":    "Persona 2: The Family Switcher\nAge: 34–45\nJob: Parent, household decision-maker\nMotivation: Safe for kids, cost-effective",
        "s3-persona3":    "Persona 3: The B2B Buyer\nAge: 30–50\nJob: Procurement / Office Manager\nMotivation: ESG reporting + bulk pricing",
        "s4-heading":     "2026 Annual Targets",
        "s4-goal1":       "Grow qualified leads by 60% YoY (from 18K to 29K MQLs)",
        "s4-goal2":       "Increase brand awareness index from 12% to 22% (Nielsen tracking)",
        "s4-goal3":       "Reduce blended CAC from $48 to $34 through owned channel growth",
        "s4-budget-total":"Marketing Budget\n$1,200,000",
        "s5-heading":     "Channel Mix",
        "s5-ch1":         "Organic Social\n28% of budget\n320K followers target",
        "s5-ch2":         "Paid Social (Meta/TikTok)\n32% of budget\n$18 target CPL",
        "s5-ch3":         "Email & CRM\n18% of budget\n42% open rate target",
        "s5-ch4":         "Influencer & Creator\n22% of budget\n85 active partners",
        "s6-heading":     "Content Pillars",
        "s6-pillar1":     "Pillar 1\nScience of Clean — ingredient education, lab stories, certifications",
        "s6-pillar2":     "Pillar 2\nReal People, Real Results — UGC, before/after, community spotlights",
        "s6-pillar3":     "Pillar 3\nSustainability Impact — carbon reports, packaging innovation, B Corp progress",
        "s7-heading":     "Campaign Calendar",
        "s7-c1":          "Q1: New Year, New Routine\nGoal: Acquisition · Channels: Paid social + email · Budget: $280K",
        "s7-c2":          "Q2: Earth Month Awareness\nGoal: Brand · Channels: PR + creator · Budget: $260K",
        "s7-c3":          "Q3: Back-to-School Bundle\nGoal: Conversion · Channels: Email + paid · Budget: $340K",
        "s7-c4":          "Q4: Gift Season\nGoal: Revenue · Channels: All channels · Budget: $320K",
        "s8-heading":     "How We Measure Success",
        "s8-m1":          "MQLs\nTarget: 29,000",
        "s8-m2":          "Blended CAC\nTarget: $34",
        "s8-m3":          "Email Revenue\nTarget: $480K",
        "s8-m4":          "Brand Awareness\nTarget: 22%",
        "s8-reporting":   "Monthly marketing dashboard reviewed by CMO + CEO. Quarterly deep-dive with board. Real-time Looker Studio shared with sales.",
        "s9-heading":     "Where the Money Goes",
        "s9-alloc1":      "Digital Ads — 32%  ($384K)",
        "s9-alloc2":      "Content Creation — 20%  ($240K)",
        "s9-alloc3":      "Influencer & Events — 26%  ($312K)",
        "s9-alloc4":      "Tools, Tech & Ops — 22%  ($264K)",
        "s9-total":       "TOTAL\n$1,200,000",
        "s10-heading":    "Team & Responsibilities",
        "s10-m1":         "Chloe Adeyemi\nCMO\nStrategy, board reporting, brand vision",
        "s10-m2":         "Finn O'Brien\nHead of Paid\nPPC, social ads, CRO",
        "s10-m3":         "Mei Lin\nContent Director\nEditorial, creator partnerships, SEO",
        "s10-m4":         "Priya Nair\nMarketing Ops\nCRM, analytics, attribution",
        "s11-heading":    "Let's Grow Together",
        "s11-contact":    "marketing@bloomliving.co  ·  @bloomliving",
    },
    "sales_presentation": {
        "s1-greeting":    "PRESENTED TO MERIDIAN FINANCIAL GROUP",
        "s1-title":       "How Apex CRM Helps You Hit 140% of Revenue Target",
        "s1-rep":         "Jordan Blake  ·  Senior Account Executive  ·  jordan@apexcrm.com",
        "s2-heading":     "What We'll Cover Today",
        "s2-item1":       "01  About Apex CRM — 5 min",
        "s2-item2":       "02  Your Current Challenges — 10 min",
        "s2-item3":       "03  Our Solution for Meridian — 10 min",
        "s2-item4":       "04  Next Steps — 5 min",
        "s3-heading":     "Why 9,400+ Teams Choose Us",
        "s3-creds1":      "9,400+\nCustomers",
        "s3-creds2":      "62 Countries",
        "s3-creds3":      "11 Years\nin Market",
        "s3-creds4":      "#1 Rated\nEnterprise CRM",
        "s3-logos":       "Trusted by Stripe, Shopify, Canva, Figma, and 200+ enterprise customers",
        "s4-heading":     "Does This Sound Familiar?",
        "s4-pain1":       "Reps spend 40% of their day on data entry instead of selling",
        "s4-pain2":       "Forecast accuracy is below 65% — leadership can't plan headcount",
        "s4-pain3":       "No single view of the customer — CRM, support, and billing are siloed",
        "s4-cost":        "Cost of Inaction\n$1.8M per year",
        "s5-heading":     "Apex CRM for Meridian",
        "s5-capability1": "AI Data Entry — auto-logs calls, emails, and meetings. Zero manual input.",
        "s5-capability2": "Revenue Intelligence — ML forecasting at 94% accuracy. Built for CFOs.",
        "s5-capability3": "360° Customer View — CRM + support + billing unified in one timeline.",
        "s6-heading":     "What Meridian Can Expect",
        "s6-roi1":        "45% reduction in time spent on data entry",
        "s6-roi2":        "$1.1M saved per year in rep productivity",
        "s6-roi3":        "6 hours saved per rep per week",
        "s6-payback":     "Payback period: 4 months. Based on 12 similar financial services customers.",
        "s7-heading":     "Results Our Customers Achieve",
        "s7-case1":       "Pinnacle Bank\nFinancial Services\n\"Forecast accuracy went from 58% to 91% in 60 days. Our CFO finally trusts the number.\"",
        "s7-case2":       "Creston Capital\nAsset Management\n\"We cut CRM admin from 8 hrs/week to under 90 minutes. Reps are happier and quota attainment is up 28%.\"",
        "s8-heading":     "Your Tailored Package",
        "s8-package":     "Meridian Growth Plan\n$420/seat/mo (85 seats)\nAI Data Entry · Revenue Intelligence · 360 View · Dedicated CSM · SLA 99.9%",
        "s8-onboarding":  "White-glove onboarding: 2-week data migration, custom field mapping, 4 live training sessions. Go-live in 30 days guaranteed.",
        "s9-heading":     "Ready to Get Started?",
        "s9-step1":       "1. Technical demo with your IT team — this week",
        "s9-step2":       "2. Pilot agreement signed — within 2 business days",
        "s9-step3":       "3. Kickoff + data migration — Week 1",
        "s9-contact":     "jordan@apexcrm.com  ·  +1 (415) 882-3300",
    },
    "training_module": {
        "s1-module":      "MODULE 03 — SALES ENABLEMENT",
        "s1-title":       "Objection Handling Mastery",
        "s1-trainer":     "Trainer: Sarah Winters  ·  Duration: 60 minutes",
        "s2-heading":     "What You'll Learn",
        "s2-obj1":        "01  Recognise the 5 most common buyer objections",
        "s2-obj2":        "02  Apply the LAER framework in real conversations",
        "s2-obj3":        "03  Turn price objections into value conversations",
        "s2-obj4":        "04  Practice with live role-play scenarios",
        "s3-heading":     "Session Overview",
        "s3-block1":      "Module 1 — Understanding Objections — 15 min",
        "s3-block2":      "Module 2 — The LAER Framework — 20 min",
        "s3-block3":      "Module 3 — Price & Competition — 15 min",
        "s3-break":       "Role-play exercises & Q&A — 10 min",
        "s4-section":     "SECTION 1: UNDERSTANDING OBJECTIONS",
        "s4-heading":     "Why Buyers Object",
        "s4-definition":  "An objection is not a rejection — it's a request for more information or reassurance. 80% of sales are made after the 5th follow-up.",
        "s4-key1":        "Objections signal engagement, not disinterest",
        "s4-key2":        "Most objections are really questions in disguise",
        "s4-key3":        "Preparation eliminates 70% of objections before they arise",
        "s5-heading":     "The LAER Framework",
        "s5-step1":       "Step 1\nListen — let them finish completely",
        "s5-step2":       "Step 2\nAcknowledge — validate their concern",
        "s5-step3":       "Step 3\nExplore — ask a clarifying question",
        "s5-step4":       "Step 4\nRespond — address with evidence",
        "s6-heading":     "Real-World Example",
        "s6-scenario":    "Scenario: Enterprise prospect says 'Your price is 30% higher than your competitor.'",
        "s6-solution":    "Solution: 'That's fair — let me show you the 3 capabilities they don't include and the ROI our last 5 customers of your size achieved.'",
        "s6-outcome":     "Outcome: Prospect requests ROI case study. Deal progresses to technical eval.",
        "s7-section":     "SECTION 2: PRICE OBJECTIONS",
        "s7-heading":     "Turning Price into Value",
        "s7-content":     "Price objections are almost always value gaps. The buyer hasn't yet connected your price to a measurable outcome they care about. Your job is to bridge that gap.",
        "s8-heading":     "Before vs. After",
        "s8-before":      "WITHOUT KNOWLEDGE\nPanics under price pressure · Discounts immediately · Loses margin · Loses confidence",
        "s8-after":       "WITH KNOWLEDGE\nRedirects to ROI · Defends value confidently · Protects margin · Closes at full price",
        "s9-heading":     "Do's and Don'ts",
        "s9-do1":         "DO: Always acknowledge before responding",
        "s9-do2":         "DO: Use customer proof points to answer objections",
        "s9-do3":         "DO: Ask 'What would need to be true for this to make sense?'",
        "s9-dont1":       "DON'T: Interrupt or argue with the objection",
        "s9-dont2":       "DON'T: Discount without exploring the real concern first",
        "s9-dont3":       "DON'T: Make promises you can't keep under pressure",
        "s10-heading":    "Test Your Understanding",
        "s10-q1":         "Q1: A prospect says 'We need to think about it.' What is the LAER first step?\nA) Immediately send a follow-up email\nB) Listen fully, then acknowledge their concern",
        "s10-q2":         "Q2: Price objections are usually really about…\nA) The budget being too low\nB) A perceived value gap that hasn't been bridged",
        "s11-heading":    "Key Takeaways",
        "s11-t1":         "Objections are buying signals — treat them as questions",
        "s11-t2":         "LAER gives you a repeatable framework for any objection",
        "s11-t3":         "Value conversations beat price discounts every time",
        "s11-next":       "Next: Module 04 — Negotiation & Closing Techniques",
        "s12-heading":    "Questions & Discussion",
        "s12-contact":    "sarah.winters@company.com  ·  Resources: company.com/training",
    },
    "executive_strategy": {},
}


def get_fallback_content(template: dict) -> dict[str, str]:
    """Return hardcoded realistic demo content for a template."""
    return FALLBACK.get(template.get("slug", ""), {})


# ─── Colour utilities ─────────────────────────────────────────────────────────

def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def is_dark(hex_color: str) -> bool:
    r, g, b = hex_to_rgb(hex_color)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255 < 0.5


def contrast_color(bg_hex: str) -> str:
    return "#FFFFFF" if is_dark(bg_hex) else DARK_TEXT_COLOR


def resolve_bg_color(bg: str, theme_colors: dict) -> str | None:
    if not bg or bg == "transparent":
        return None
    if bg in SEED_DARK_MAP:
        return theme_colors.get(SEED_DARK_MAP[bg], bg)
    if bg.upper() in {c.upper() for c in NEUTRAL_BACKGROUNDS}:
        return theme_colors.get("surface", bg)
    return bg


def resolve_text_color(color: str, bg: str, theme_colors: dict) -> str:
    if color and color != "transparent":
        return color
    resolved_bg = resolve_bg_color(bg, theme_colors)
    if resolved_bg:
        return contrast_color(resolved_bg)
    return theme_colors.get("text", "#000000")


# ─── Loaders ─────────────────────────────────────────────────────────────────

def load_themes() -> dict[str, dict]:
    themes = {}
    for path in THEMES_DIR.glob("*.json"):
        data = json.loads(path.read_text(encoding="utf-8"))
        themes[data["name"]] = data
    return themes


def load_templates() -> list[dict]:
    templates = []
    for path in sorted(TEMPLATES_DIR.glob("*.json")):
        data = json.loads(path.read_text(encoding="utf-8"))
        data["slug"] = path.stem
        templates.append(data)
    return templates


# ─── PPTX builder ─────────────────────────────────────────────────────────────

def _px_to_emu(px: float, design_dim: float, slide_dim_in: float) -> int:
    return int((px / design_dim) * slide_dim_in * EMU_PER_INCH)


def _font_family_for_type(block_type: str, fonts: dict) -> str:
    if block_type in ("title", "heading", "subtitle"):
        return fonts.get("heading", {}).get("family", "Calibri")
    if block_type == "caption":
        return fonts.get("caption", {}).get("family", "Calibri")
    return fonts.get("body", {}).get("family", "Calibri")


def _logo_png_bytes(svg_path: Path, height_px: int = 80) -> bytes | None:
    """Convert logo SVG to PNG bytes at the requested height. Returns None if unavailable."""
    try:
        import cairosvg
        return cairosvg.svg2png(url=str(svg_path), output_height=height_px)
    except Exception as exc:
        print(f"    [WARN] Logo conversion failed: {exc}")
        return None


def build_pptx(slides: list[dict], theme: dict, content: dict, logo_path: Path | None = None):
    """
    Build and return a python-pptx Presentation with Gamma-quality theming.

    Args:
        slides:     List of slide dicts from the template JSON.
        theme:      Theme dict with 'colors' and 'fonts' keys.
        content:    Dict mapping block_id → replacement text.
        logo_path:  Optional path to logo SVG — inserted top-left on slide 1 only.
    """
    from pptx import Presentation as PptxPresentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    ALIGN_MAP = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }

    colors     = theme["colors"]
    fonts      = theme["fonts"]
    bg_hex     = colors.get("background", "#FFFFFF")
    accent_hex = colors.get("accent", "#000000")

    prs = PptxPresentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]

    # Pre-convert logo once
    logo_bytes = _logo_png_bytes(logo_path) if logo_path and logo_path.exists() else None

    def rgb(hex_color: str) -> RGBColor:
        r, g, b = hex_to_rgb(hex_color)
        return RGBColor(r, g, b)

    def no_line(shape) -> None:
        shape.line.fill.background()

    sorted_slides = sorted(slides, key=lambda s: s.get("order", 0))

    for slide_idx, slide_data in enumerate(sorted_slides):
        slide = prs.slides.add_slide(blank_layout)

        # ── Full-bleed background ──────────────────────────────────────────
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(bg_hex)

        # ── Top accent stripe ──────────────────────────────────────────────
        stripe_h = _px_to_emu(6, DESIGN_H, SLIDE_H_IN)
        stripe = slide.shapes.add_shape(
            1, 0, 0, int(SLIDE_W_IN * EMU_PER_INCH), stripe_h
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = rgb(accent_hex)
        no_line(stripe)

        # ── Logo on first slide only ───────────────────────────────────────
        if slide_idx == 0 and logo_bytes:
            logo_h_in  = 0.6          # ~58px at 96dpi — compact, top-left
            logo_left  = Inches(0.35)
            logo_top   = Inches(0.12)
            logo_stream = io.BytesIO(logo_bytes)
            slide.shapes.add_picture(logo_stream, logo_left, logo_top, height=Inches(logo_h_in))

        # ── Content blocks ────────────────────────────────────────────────
        for block in slide_data.get("blocks", []):
            if block.get("type") == "image":
                continue

            pos        = block.get("position", {"x": 0, "y": 0, "w": 200, "h": 100})
            styling    = block.get("styling", {})
            block_type = block.get("type", "text")
            block_id   = block.get("id", "")

            left   = _px_to_emu(pos.get("x", 0),   DESIGN_W, SLIDE_W_IN)
            top    = _px_to_emu(pos.get("y", 0),   DESIGN_H, SLIDE_H_IN)
            width  = _px_to_emu(pos.get("w", 200), DESIGN_W, SLIDE_W_IN)
            height = _px_to_emu(pos.get("h", 100), DESIGN_H, SLIDE_H_IN)

            raw_text = content.get(block_id) or block.get("content", "")
            text     = _re.sub(r'\[PLACEHOLDER:[^\]]*\]|\[Placeholder:[^\]]*\]', '', raw_text).strip() or raw_text

            raw_bg      = styling.get("background_color", "transparent")
            resolved_bg = resolve_bg_color(raw_bg, colors)
            text_color  = resolve_text_color(styling.get("color", "transparent"), raw_bg, colors)
            font_name   = _font_family_for_type(block_type, fonts)
            font_size   = Pt(styling.get("font_size", 20))
            bold        = styling.get("font_weight", 400) >= 700
            alignment   = ALIGN_MAP.get(styling.get("text_align", "left"), PP_ALIGN.LEFT)

            if resolved_bg:
                # Gamma-style card: rounded rectangle + internal padding
                shape = slide.shapes.add_shape(5, left, top, width, height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(resolved_bg)
                no_line(shape)

                pad = _px_to_emu(CARD_PADDING_PX, DESIGN_W, SLIDE_W_IN)
                tf  = shape.text_frame
                tf.word_wrap    = True
                tf.margin_left  = pad
                tf.margin_right = pad
                tf.margin_top   = pad
                tf.margin_bottom = pad

                p   = tf.paragraphs[0]
                p.alignment = ALIGN_MAP.get(styling.get("text_align", "center"), PP_ALIGN.CENTER)
                run = p.add_run()
                run.text           = text
                run.font.size      = font_size
                run.font.bold      = bold
                run.font.color.rgb = rgb(text_color)
                run.font.name      = font_name
            else:
                # Plain textbox
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf    = txBox.text_frame
                tf.word_wrap = True

                p   = tf.paragraphs[0]
                p.alignment = alignment
                run = p.add_run()
                run.text           = text
                run.font.size      = font_size
                run.font.bold      = bold
                run.font.color.rgb = rgb(text_color)
                run.font.name      = font_name

    return prs


# ─── Gemini content generation ────────────────────────────────────────────────

def get_content(template: dict) -> dict[str, str]:
    """Return demo content for a template (hardcoded fallback, no Gemini)."""
    return get_fallback_content(template)


# ─── Orchestration ────────────────────────────────────────────────────────────

async def process_template(template: dict, themes: dict, force: bool, db_session) -> bool:
    from sqlalchemy import text as sql_text

    slug       = template["slug"]
    name       = template["name"]
    theme_name = template.get("theme_name", "")
    theme      = themes.get(theme_name)

    if not theme:
        print(f"  [{name}] SKIP — theme '{theme_name}' not found in seeds/themes/")
        return False

    out_path = PREVIEWS_DIR / f"{slug}.pptx"
    if out_path.exists() and not force:
        print(f"  [{name}] SKIP — {out_path.name} already exists (use --force to regenerate)")
        return False

    content = get_content(template)
    print(f"  [{name}] building PPTX (theme: {theme_name}, {len(content)} blocks)...")
    prs = build_pptx(template["slides"], theme, content, logo_path=LOGO_PATH)

    PREVIEWS_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"  [{name}] saved -> {out_path.name}")

    preview_url = f"/previews/{slug}.pptx"
    await db_session.execute(
        sql_text("UPDATE templates SET preview_pptx_path = :url WHERE name = :name"),
        {"url": preview_url, "name": name},
    )
    await db_session.commit()
    print(f"  [{name}] DB updated OK")
    return True


async def main(force: bool = False, only_slug: str | None = None) -> None:
    import app.core.database_models  # noqa: F401
    from app.core.database import close_db, init_db
    import app.core.database as _db_module

    await init_db()

    themes    = load_themes()
    templates = load_templates()

    if only_slug:
        templates = [t for t in templates if t["slug"] == only_slug]
        if not templates:
            print(f"ERROR: no template with slug '{only_slug}' found in {TEMPLATES_DIR}")
            await close_db()
            return

    print(f"\nGenerating previews for {len(templates)} template(s)...\n")
    ok = skipped = 0

    async with _db_module._session_factory() as db:
        for i, template in enumerate(templates):
            success = await process_template(template, themes, force, db)
            if success:
                ok += 1
            else:
                skipped += 1

    await close_db()
    print(f"\nDone. {ok} generated, {skipped} skipped.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate preview PPTXs for all templates")
    parser.add_argument("--force",    action="store_true", help="Regenerate even if file exists")
    parser.add_argument("--template", metavar="SLUG",      help="Single template slug (e.g. business_pitch)")
    args = parser.parse_args()
    asyncio.run(main(force=args.force, only_slug=args.template))
