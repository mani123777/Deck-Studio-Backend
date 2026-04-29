import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from seeds.generate_preview_pptx import hex_to_rgb, is_dark, contrast_color


def test_hex_to_rgb_6digit():
    r, g, b = hex_to_rgb("#1A1A2E")
    assert r == 26 and g == 26 and b == 46


def test_hex_to_rgb_strips_hash():
    r, g, b = hex_to_rgb("FFFFFF")
    assert r == 255 and g == 255 and b == 255


def test_is_dark_dark_color():
    assert is_dark("#1A1A2E") is True


def test_is_dark_light_color():
    assert is_dark("#FFFFFF") is False


def test_contrast_color_dark_bg():
    assert contrast_color("#1A1A2E") == "#FFFFFF"


def test_contrast_color_light_bg():
    assert contrast_color("#FFFFFF") == "#1A1A2E"


from seeds.generate_preview_pptx import load_themes, load_templates


def test_load_themes_returns_dict_keyed_by_name():
    themes = load_themes()
    assert "Minimal Professional" in themes
    assert "Startup Energy" in themes
    assert "colors" in themes["Minimal Professional"]
    assert "fonts" in themes["Minimal Professional"]


def test_load_templates_returns_all_nine():
    templates = load_templates()
    slugs = [t["slug"] for t in templates]
    assert "business_pitch" in slugs
    assert len(templates) == 9


def test_template_has_theme_name():
    templates = load_templates()
    for t in templates:
        assert "theme_name" in t, f"{t['slug']} missing theme_name"


def test_template_theme_exists_in_themes():
    themes = load_themes()
    templates = load_templates()
    for t in templates:
        assert t["theme_name"] in themes, \
            f"{t['slug']} theme '{t['theme_name']}' not found in seeds/themes/"


from seeds.generate_preview_pptx import resolve_bg_color, resolve_text_color

THEME_COLORS = {
    "primary": "#6C63FF",
    "secondary": "#FF6584",
    "accent": "#43D9AD",
    "background": "#FAFAFA",
    "text": "#2D2D2D",
    "surface": "#EFEFEF",
}


def test_resolve_bg_seed_dark_blue_maps_to_primary():
    assert resolve_bg_color("#1A1A2E", THEME_COLORS) == "#6C63FF"


def test_resolve_bg_seed_secondary_maps_to_secondary():
    assert resolve_bg_color("#16213E", THEME_COLORS) == "#FF6584"


def test_resolve_bg_seed_accent_maps_to_accent():
    assert resolve_bg_color("#0F3460", THEME_COLORS) == "#43D9AD"


def test_resolve_bg_neutral_maps_to_surface():
    assert resolve_bg_color("#F5F5F5", THEME_COLORS) == "#EFEFEF"


def test_resolve_bg_transparent_returns_none():
    assert resolve_bg_color("transparent", THEME_COLORS) is None


def test_resolve_bg_custom_color_passthrough():
    assert resolve_bg_color("#ABCDEF", THEME_COLORS) == "#ABCDEF"


def test_resolve_text_color_uses_block_color():
    assert resolve_text_color("#FF0000", "transparent", THEME_COLORS) == "#FF0000"


def test_resolve_text_color_on_dark_bg_returns_white():
    assert resolve_text_color("transparent", "#6C63FF", THEME_COLORS) == "#FFFFFF"


def test_resolve_text_color_fallback_to_theme_text():
    assert resolve_text_color("transparent", "transparent", THEME_COLORS) == "#2D2D2D"


from seeds.generate_preview_pptx import build_pptx

MINIMAL_THEME = {
    "colors": {
        "primary": "#1A1A2E",
        "secondary": "#16213E",
        "accent": "#0F3460",
        "background": "#FFFFFF",
        "text": "#1A1A2E",
        "surface": "#F5F5F5",
    },
    "fonts": {
        "heading": {"family": "Inter", "size": 56, "weight": 700},
        "body": {"family": "Inter", "size": 28, "weight": 400},
        "caption": {"family": "Inter", "size": 20, "weight": 300},
    },
}

MINIMAL_SLIDES = [
    {
        "order": 1,
        "type": "title",
        "blocks": [
            {
                "id": "s1-title",
                "type": "title",
                "content": "Test Title",
                "position": {"x": 160, "y": 380, "w": 1600, "h": 150},
                "styling": {
                    "font_family": "Inter",
                    "font_size": 80,
                    "font_weight": 700,
                    "color": "#1A1A2E",
                    "background_color": "transparent",
                    "text_align": "center",
                },
            }
        ],
    }
]


def test_build_pptx_returns_presentation_object():
    from pptx.presentation import Presentation as PptxPresentation
    prs = build_pptx(MINIMAL_SLIDES, MINIMAL_THEME, {})
    assert isinstance(prs, PptxPresentation)


def test_build_pptx_correct_slide_count():
    prs = build_pptx(MINIMAL_SLIDES, MINIMAL_THEME, {})
    assert len(prs.slides) == 1


def test_build_pptx_slide_dimensions():
    from pptx.util import Inches
    prs = build_pptx(MINIMAL_SLIDES, MINIMAL_THEME, {})
    assert abs(prs.slide_width - Inches(13.33)) < 1000
    assert abs(prs.slide_height - Inches(7.5)) < 1000


def test_build_pptx_content_override():
    content = {"s1-title": "Overridden Title"}
    prs = build_pptx(MINIMAL_SLIDES, MINIMAL_THEME, content)
    slide = prs.slides[0]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert any("Overridden Title" in t for t in texts)
