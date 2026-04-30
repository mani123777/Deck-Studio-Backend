from __future__ import annotations

import io
import uuid
from pathlib import Path
from typing import Any

from fastapi import APIRouter, Depends, File, HTTPException, Request, UploadFile
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.api.dependencies import get_current_user
from app.core.database import get_db
from app.models.theme import Theme
from app.models.user import User
from app.schemas.presentation import (
    BlockSchema,
    CreatePresentationRequest,
    PositionSchema,
    SlideBackgroundSchema,
    SlideSchema,
    StylingSchema,
)
from app.services import presentation_service
from app.utils.logger import get_logger

router = APIRouter(prefix="/import", tags=["import"])
logger = get_logger(__name__)

# Editor canvas dimensions (same as slide_generator_agent.py W, H)
CANVAS_W = 1280
CANVAS_H = 720

# Where extracted slide images are stored / served
_BACKEND_DIR = Path(__file__).resolve().parent.parent.parent.parent  # → backend/
IMPORTS_DIR = _BACKEND_DIR / "storage" / "imports"


# ---------------------------------------------------------------------------
# Coordinate helpers
# ---------------------------------------------------------------------------

def _emu_to_canvas(emu: int | None, slide_dim_emu: int, canvas_px: int) -> float:
    """Map an EMU coordinate to the 1280×720 canvas pixel space."""
    if not emu or not slide_dim_emu:
        return 0.0
    return round((emu / slide_dim_emu) * canvas_px, 1)


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def _rgb_to_hex(rgb: Any) -> str | None:
    """Convert python-pptx RGBColor / tuple to '#rrggbb'. Returns None on failure."""
    try:
        return f"#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}"
    except Exception:
        return None


def _fill_to_hex(fill: Any) -> str | None:
    """Return a hex color string from a pptx FillFormat, or None if not solid."""
    try:
        from pptx.enum.dml import MSO_THEME_COLOR
        if fill.type is None:
            return None
        from pptx.enum.dml import MSO_FILL_TYPE  # type: ignore
        if fill.type.name != "SOLID":
            return None
        return _rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Text alignment
# ---------------------------------------------------------------------------

_ALIGN_MAP: dict[Any, str] = {}

def _pptx_align_to_str(align: Any) -> str:
    try:
        from pptx.enum.text import PP_ALIGN
        return {
            PP_ALIGN.LEFT:    "left",
            PP_ALIGN.CENTER:  "center",
            PP_ALIGN.RIGHT:   "right",
            PP_ALIGN.JUSTIFY: "justify",
        }.get(align, "left")
    except Exception:
        return "left"


# ---------------------------------------------------------------------------
# Shape → BlockSchema
# ---------------------------------------------------------------------------

def _shape_to_block(
    shape: Any,
    slide_w_emu: int,
    slide_h_emu: int,
    folder_id: str,
    base_url: str,
    image_counter: list[int],
) -> BlockSchema | None:
    """Convert one python-pptx shape to a BlockSchema, or return None to skip."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # --- Position -----------------------------------------------------------
    x = _emu_to_canvas(shape.left,  slide_w_emu, CANVAS_W)
    y = _emu_to_canvas(shape.top,   slide_h_emu, CANVAS_H)
    w = _emu_to_canvas(shape.width, slide_w_emu, CANVAS_W)
    h = _emu_to_canvas(shape.height,slide_h_emu, CANVAS_H)

    # Clamp: never allow negative size
    w = max(w, 10.0)
    h = max(h, 10.0)

    pos = PositionSchema(x=x, y=y, w=w, h=h)

    # --- Picture ------------------------------------------------------------
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img_data = shape.image.blob
            img_ext  = shape.image.ext or "png"
            idx      = image_counter[0]
            image_counter[0] += 1

            dest_dir = IMPORTS_DIR / folder_id
            dest_dir.mkdir(parents=True, exist_ok=True)
            img_path = dest_dir / f"img_{idx}.{img_ext}"
            img_path.write_bytes(img_data)

            img_url = f"{base_url}/imports/{folder_id}/img_{idx}.{img_ext}"
            return BlockSchema(
                id=str(uuid.uuid4()),
                type="image",
                content=img_url,
                position=pos,
                styling=StylingSchema(),
            )
        except Exception as exc:
            logger.warning(f"Skipping picture shape: {exc}")
            return None

    # --- Text shapes --------------------------------------------------------
    if not shape.has_text_frame:
        # Non-text, non-picture shape — try to render as a colored shape block
        try:
            fill_hex = _fill_to_hex(shape.fill)
            if fill_hex:
                return BlockSchema(
                    id=str(uuid.uuid4()),
                    type="shape",
                    content="",
                    position=pos,
                    styling=StylingSchema(background_color=fill_hex, color=fill_hex),
                )
        except Exception:
            pass
        return None

    tf = shape.text_frame

    # Collect paragraphs, preserving multi-line content
    lines: list[str] = [p.text for p in tf.paragraphs if p.text.strip()]
    text = "\n".join(lines)
    if not text:
        return None

    # --- Detect title placeholder -------------------------------------------
    is_title = False
    try:
        ph = shape.placeholder_format
        if ph is not None:
            is_title = ph.idx in (0, 1)
    except Exception:
        pass

    # --- Font properties: read from first non-empty run ----------------------
    font_size   = 36 if is_title else 18
    font_weight = 700 if is_title else 400
    color       = "#1f2937"
    text_align  = "left"

    try:
        for para in tf.paragraphs:
            # Paragraph-level alignment takes precedence
            if para.alignment is not None:
                text_align = _pptx_align_to_str(para.alignment)

            for run in para.runs:
                f = run.font
                try:
                    if f.size:
                        font_size = max(8, min(96, int(f.size.pt)))
                except Exception:
                    pass
                try:
                    if f.bold:
                        font_weight = 700
                except Exception:
                    pass
                try:
                    if f.color and f.color.type is not None:
                        hex_val = _rgb_to_hex(f.color.rgb)
                        if hex_val:
                            color = hex_val
                except Exception:
                    pass
                break  # only first run per paragraph
            if lines:
                break   # only first non-empty paragraph
    except Exception:
        pass

    return BlockSchema(
        id=str(uuid.uuid4()),
        type="text",
        content=text,
        position=pos,
        styling=StylingSchema(
            font_size=font_size,
            font_weight=font_weight,
            color=color,
            text_align=text_align,
        ),
    )


# ---------------------------------------------------------------------------
# Slide background
# ---------------------------------------------------------------------------

def _slide_background(slide: Any) -> SlideBackgroundSchema:
    """Extract the slide's background fill color."""
    try:
        bg_fill = slide.background.fill
        hex_col = _fill_to_hex(bg_fill)
        if hex_col:
            return SlideBackgroundSchema(type="color", value=hex_col)
    except Exception:
        pass
    return SlideBackgroundSchema(type="color", value="#ffffff")


# ---------------------------------------------------------------------------
# Full PPTX → slides parser
# ---------------------------------------------------------------------------

def _parse_pptx(data: bytes, folder_id: str, base_url: str) -> tuple[str, list[SlideSchema]]:
    """Parse PPTX bytes into editable (title, slides) using the same
    1280×720 coordinate system as the slide generator."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))

    slide_w_emu = int(prs.slide_width)  if prs.slide_width  else 12_192_000
    slide_h_emu = int(prs.slide_height) if prs.slide_height else  6_858_000

    presentation_title = "Imported Presentation"
    slides: list[SlideSchema] = []
    image_counter = [0]  # mutable counter shared across slides

    for slide_idx, slide in enumerate(prs.slides):
        blocks: list[BlockSchema] = []
        background = _slide_background(slide)

        for shape in slide.shapes:
            try:
                block = _shape_to_block(
                    shape, slide_w_emu, slide_h_emu,
                    folder_id, base_url, image_counter,
                )
            except Exception as exc:
                logger.warning(f"Slide {slide_idx} shape skipped: {exc}")
                continue

            if block is None:
                continue

            # Grab title from first title-placeholder on first slide
            if slide_idx == 0 and block.type == "text" and block.styling.font_size >= 28:
                first_line = block.content.split("\n")[0].strip()
                if first_line and presentation_title == "Imported Presentation":
                    presentation_title = first_line[:120]

            blocks.append(block)

        # Sort blocks top-to-bottom so reading order matches visual order
        blocks.sort(key=lambda b: (b.position.y, b.position.x))

        if not blocks:
            blocks.append(BlockSchema(
                id=str(uuid.uuid4()),
                type="text",
                content=f"Slide {slide_idx + 1}",
                position=PositionSchema(x=64, y=280, w=CANVAS_W - 128, h=80),
                styling=StylingSchema(font_size=32, font_weight=700, color="#1f2937", text_align="center"),
            ))

        has_large_text = any(b.type == "text" and b.styling.font_size >= 30 for b in blocks)
        slide_type = "title" if (slide_idx == 0 and has_large_text) else "content"

        slides.append(SlideSchema(
            order=slide_idx,
            type=slide_type,
            background=background,
            blocks=blocks,
        ))

    return presentation_title, slides


# ---------------------------------------------------------------------------
# Endpoint
# ---------------------------------------------------------------------------

@router.post("/pptx")
async def import_pptx(
    request: Request,
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
) -> Any:
    """Import a .pptx file: converts every shape to an editable block using
    the same 1280×720 coordinate system as the slide generator."""
    filename = (file.filename or "").lower()
    if not filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")

    data = await file.read()
    if len(data) > 50 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large (max 50 MB)")

    folder_id = str(uuid.uuid4())
    base_url  = str(request.base_url).rstrip("/")

    try:
        title, slides = _parse_pptx(data, folder_id, base_url)
    except Exception as exc:
        logger.error(f"PPTX parse error: {exc}")
        raise HTTPException(status_code=422, detail=f"Could not parse PPTX file: {exc}")

    theme = (await db.execute(select(Theme))).scalars().first()
    if not theme:
        raise HTTPException(status_code=500, detail="No theme found in database")

    req = CreatePresentationRequest(
        title=title,
        description="Imported from PowerPoint",
        slides=slides,
        theme_id=str(theme.id),
        template_id="",
        logo_url="",
    )

    presentation = await presentation_service.create_presentation(db, current_user, req)
    logger.info(
        f"Imported PPTX '{title}' ({len(slides)} slides, {folder_id}) "
        f"→ {presentation.id} for user {current_user.id}"
    )
    return {"id": presentation.id, "title": presentation.title, "total_slides": presentation.total_slides}
