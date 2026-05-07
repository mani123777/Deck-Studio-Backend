from __future__ import annotations

import io
import re
from pathlib import Path
from urllib.parse import urlparse

from app.core.storage import export_path
from app.models.presentation import Presentation
from app.models.theme import Theme
from app.utils.logger import get_logger

logger = get_logger(__name__)

# Must match the editor canvas (slide_generator_agent.py W, H)
CANVAS_W = 1280
CANVAS_H = 720
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Backend root — used to resolve local /imports/ and /previews/ URLs
_BACKEND_DIR = Path(__file__).resolve().parent.parent.parent.parent  # → backend/


def _px_to_emu(px: float, canvas_dim: float, slide_dim_in: float) -> int:
    return int((px / canvas_dim) * slide_dim_in * 914_400)


def _hex_to_rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _extract_color(value: str) -> str | None:
    """Return a usable hex color from any CSS color value or gradient string.
    Returns None if no color can be extracted."""
    if not value or value in ("transparent", "none", ""):
        return None
    value = value.strip()
    # Plain hex
    if re.match(r"^#[0-9a-fA-F]{3,6}$", value):
        return value
    # Hex inside a gradient / other string
    hexes = re.findall(r"#[0-9a-fA-F]{6}|#[0-9a-fA-F]{3}", value)
    if hexes:
        return hexes[0]
    # rgba(r, g, b, ...)
    m = re.search(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", value)
    if m:
        return f"#{int(m.group(1)):02x}{int(m.group(2)):02x}{int(m.group(3)):02x}"
    return None


def _resolve_image(url: str) -> io.BytesIO | None:
    """Resolve an image URL to bytes.
    Tries filesystem lookup for known static routes, then HTTP fetch."""
    if not url:
        return None
    try:
        parsed = urlparse(url)
        # parsed.path is the path component regardless of whether url is absolute or relative
        path = parsed.path  # e.g. /imports/uuid/img_0.png

        # Walk up to find the actual /imports/ segment in case there's a prefix
        # (e.g. /api/v1/imports/ or the path already starts with /imports/)
        imports_idx = path.find("/imports/")
        previews_idx = path.find("/previews/")
        storage_idx = path.find("/storage/")

        local: Path | None = None
        if imports_idx >= 0:
            rel = path[imports_idx:].lstrip("/")   # imports/uuid/img_0.png
            local = _BACKEND_DIR / "storage" / rel
        elif previews_idx >= 0:
            rel = path[previews_idx:].lstrip("/")  # previews/...
            local = _BACKEND_DIR.parent / "seeds" / rel
        elif storage_idx >= 0:
            rel = path[storage_idx + len("/storage/"):]  # strip leading /storage/
            local = _BACKEND_DIR / "storage" / rel

        if local is not None:
            if local.exists():
                return io.BytesIO(local.read_bytes())
            logger.warning(f"Image not found on disk: {local} (url={url})")

        # Fallback: HTTP fetch (CDN images, absolute external URLs)
        if parsed.scheme in ("http", "https"):
            import httpx
            resp = httpx.get(url, timeout=15, follow_redirects=True)
            if resp.status_code == 200:
                return io.BytesIO(resp.content)
            logger.warning(f"HTTP fetch for image returned {resp.status_code}: {url}")
    except Exception as exc:
        logger.warning(f"Could not resolve image {url}: {exc}")
    return None


class PptxExportAgent:
    async def run(self, presentation: Presentation, theme: Theme) -> Path:
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = PptxPresentation()
        prs.slide_width  = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)

        blank_layout = prs.slide_layouts[6]

        colors = theme.colors
        fonts  = theme.fonts

        align_map = {
            "left":    PP_ALIGN.LEFT,
            "center":  PP_ALIGN.CENTER,
            "right":   PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }

        for slide_data in sorted(presentation.slides or [], key=lambda s: s.get("order", 0)):
            pptx_slide = prs.slides.add_slide(blank_layout)

            # Presenter notes — render into the slide's notes_slide so they
            # appear in PowerPoint's notes pane and read back from
            # python-pptx-style consumers.
            notes_text = (slide_data.get("notes") or "").strip()
            if notes_text:
                try:
                    pptx_slide.notes_slide.notes_text_frame.text = notes_text
                except Exception:
                    pass

            # Aggressively remove ALL placeholder shapes (p:sp with a p:ph child).
            # pptx_slide.placeholders only enumerates shapes the library recognises;
            # some slide-master-inherited placeholders are invisible to that API but
            # still render as "[No Title]" / "[Content]" boxes in the output file.
            from pptx.oxml.ns import qn as _qn
            _sp_tree = pptx_slide.shapes._spTree
            for _sp in list(_sp_tree.findall(_qn('p:sp'))):
                if _sp.find('.//' + _qn('p:ph')) is not None:
                    _sp_tree.remove(_sp)

            # ── Slide background ──────────────────────────────────────────
            bg      = slide_data.get("background") or {}
            bg_type = bg.get("type", "color")
            bg_val  = bg.get("value", colors.get("background", "#ffffff"))
            bg_hex  = _extract_color(bg_val if bg_type != "image" else "")
            if bg_hex:
                fill = pptx_slide.background.fill
                fill.solid()
                try:
                    fill.fore_color.rgb = _hex_to_rgb(bg_hex)
                except Exception:
                    pass

            # Render in correct z-order: panels → shapes → images → text
            # (matches the browser stacking context in the editor)
            all_blocks = slide_data.get("blocks", [])
            ordered_blocks = (
                [b for b in all_blocks if b.get("type") in ("panel",)] +
                [b for b in all_blocks if b.get("type") in ("shape",)] +
                [b for b in all_blocks if b.get("type") == "image"] +
                [b for b in all_blocks if b.get("type") not in ("panel", "shape", "image")]
            )

            for block in ordered_blocks:
                btype   = block.get("type", "text")
                pos     = block.get("position", {"x": 0, "y": 0, "w": 100, "h": 100})
                styling = block.get("styling", {})
                content = block.get("content", "")

                left   = _px_to_emu(pos.get("x", 0),   CANVAS_W, SLIDE_W_IN)
                top    = _px_to_emu(pos.get("y", 0),   CANVAS_H, SLIDE_H_IN)
                width  = _px_to_emu(pos.get("w", 100), CANVAS_W, SLIDE_W_IN)
                height = _px_to_emu(pos.get("h", 100), CANVAS_H, SLIDE_H_IN)

                # ── Image block ───────────────────────────────────────────
                if btype == "image":
                    img_stream = _resolve_image(content)
                    if img_stream:
                        try:
                            # add_picture accepts a file-like object directly
                            pptx_slide.shapes.add_picture(img_stream, left, top, width, height)
                        except Exception as exc:
                            logger.warning(f"Could not embed image {content}: {exc}")
                    continue

                # ── Shape / panel → solid colored rectangle ───────────────
                if btype in ("shape", "panel"):
                    bg_color = styling.get("background_color") or styling.get("color", "")
                    hex_col  = _extract_color(bg_color)
                    if not hex_col:
                        continue
                    try:
                        shape = pptx_slide.shapes.add_shape(1, left, top, width, height)
                        shape.line.width = 0
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = _hex_to_rgb(hex_col)
                    except Exception as exc:
                        logger.warning(f"Shape block render failed: {exc}")
                    continue

                # ── All text-bearing blocks ───────────────────────────────
                font_size   = styling.get("font_size", 16)
                font_weight = styling.get("font_weight", 400)
                color_hex   = styling.get("color", "#000000")
                text_align  = styling.get("text_align", "left")
                font_family = styling.get("font_family", "")
                bg_fill     = _extract_color(styling.get("background_color", ""))

                # Editor font_size is in CSS px; PPT uses pt. 1px = 0.75pt @ 96 DPI.
                # Without this conversion a 96px hero title renders as 96pt (~128px)
                # and explodes out of its box.
                PX_TO_PT = 0.75
                font_size_pt = float(font_size) * PX_TO_PT

                # Pre-compute a safe font size so text never overflows the box.
                # All math in px to mirror the editor's CSS layout exactly.
                box_w_px = max(pos.get("w", 100), 10)
                box_h_px = max(pos.get("h", 100), 10)
                raw_lines = (content or "").split("\n")
                is_bold   = font_weight >= 700

                def _fits(fs_px: float) -> bool:
                    # Char width: ~0.58 for regular, ~0.62 for bold proportional fonts.
                    char_w = max(fs_px * (0.62 if is_bold else 0.58), 1.0)
                    line_h = fs_px * 1.4
                    chars_per_line = max(int(box_w_px / char_w), 1)
                    total_lines = 0
                    for ln in raw_lines:
                        if not ln:
                            total_lines += 1
                            continue
                        total_lines += max(1, -(-len(ln) // chars_per_line))
                    # 95% safety margin so descenders don't clip
                    return (total_lines * line_h) <= (box_h_px * 0.95)

                fitted_px = float(font_size)
                while fitted_px > 8 and not _fits(fitted_px):
                    fitted_px -= 1
                fitted_px = max(fitted_px, 8.0)
                fitted_size = fitted_px * PX_TO_PT

                txBox = pptx_slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                # Zero internal padding so text bounds match the editor box exactly
                from pptx.util import Emu
                tf.margin_left   = Emu(0)
                tf.margin_right  = Emu(0)
                tf.margin_top    = Emu(0)
                tf.margin_bottom = Emu(0)

                from pptx.oxml.ns import qn
                tf._txBody.find(qn('a:bodyPr')).set('anchor', 't')
                # We've manually fitted the font — disable PPT's auto-fit so it
                # doesn't fight our sizing.
                from pptx.enum.text import MSO_AUTO_SIZE
                tf.auto_size = MSO_AUTO_SIZE.NONE

                # Background fill on the text box
                if bg_fill:
                    try:
                        txBox.fill.solid()
                        txBox.fill.fore_color.rgb = _hex_to_rgb(bg_fill)
                    except Exception:
                        pass

                lines = raw_lines if content else [""]
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.alignment = align_map.get(text_align, PP_ALIGN.LEFT)
                    # Match editor line-height: 1.4
                    p.line_spacing = 1.4
                    p.space_before = Pt(0)
                    p.space_after  = Pt(0)
                    run = p.add_run()
                    run.text = line
                    f = run.font
                    f.size = Pt(fitted_size)
                    f.bold = font_weight >= 700
                    if font_family:
                        f.name = font_family.split(",")[0].strip().strip("'\"")
                    try:
                        f.color.rgb = _hex_to_rgb(color_hex)
                    except Exception:
                        pass

        out_path = export_path(str(presentation.id), "pptx")
        prs.save(str(out_path))
        logger.info(f"PPTX export saved to {out_path}")
        return out_path
